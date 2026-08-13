#!/usr/bin/env python3
"""Build the VYTRA Brand Identity book (HTML) and slide deck (PPTX)."""

from __future__ import annotations

import base64
import io
import os
from pathlib import Path

from PIL import Image, ImageChops, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import nsmap
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Inches, Pt

ROOT = Path("/home/user/VYTRA_Brand_Identity")
ASSETS = ROOT / "assets"
EXPORTS = ROOT / "exports"

FOREST = RGBColor(0x0E, 0x2A, 0x1C)
FOREST_DEEP = RGBColor(0x08, 0x1C, 0x12)
PULSE = RGBColor(0x6C, 0xA5, 0x32)
PULSE_DEEP = RGBColor(0x4F, 0x7E, 0x24)
SOFT = RGBColor(0xE8, 0xF3, 0xDC)
CREAM = RGBColor(0xF6, 0xF4, 0xEE)
PAPER = RGBColor(0xFF, 0xFE, 0xFB)
INK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x47, 0x54, 0x67)
LINE = RGBColor(0xE4, 0xE7, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WARM = RGBColor(0xF3, 0xF0, 0xE7)


def trim_whitespace(im: Image.Image, threshold: int = 246, pad: int = 24) -> Image.Image:
    """Crop near-white margins, keep a small pad."""
    rgb = im.convert("RGB")
    # Treat off-white compression haze as background.
    pixels = rgb.load()
    mask = Image.new("L", rgb.size, 0)
    mp = mask.load()
    w, h = rgb.size
    for y in range(h):
        for x in range(w):
            r, g, b = pixels[x, y]
            if r < threshold or g < threshold or b < threshold:
                mp[x, y] = 255
    bbox = mask.getbbox()
    if not bbox:
        return im
    l, t, r, b = bbox
    l = max(0, l - pad)
    t = max(0, t - pad)
    r = min(im.width, r + pad)
    b = min(im.height, b + pad)
    return im.crop((l, t, r, b))


def save_display(src: Path, dest: Path, max_side: int = 1400, pad: int = 28) -> Path:
    im = Image.open(src).convert("RGBA")
    # flatten on white for lockups
    bg = Image.new("RGBA", im.size, (255, 255, 255, 255))
    flat = Image.alpha_composite(bg, im).convert("RGB")
    flat = trim_whitespace(flat, pad=pad)
    flat.thumbnail((max_side, max_side), Image.Resampling.LANCZOS)
    dest.parent.mkdir(parents=True, exist_ok=True)
    flat.save(dest, "PNG", optimize=True)
    return dest


def encode_png(path: Path, max_side: int = 720, pad: int = 20) -> str:
    im = Image.open(path).convert("RGBA")
    bg = Image.new("RGBA", im.size, (255, 255, 255, 255))
    flat = Image.alpha_composite(bg, im).convert("RGB")
    flat = trim_whitespace(flat, pad=pad)
    flat.thumbnail((max_side, max_side), Image.Resampling.LANCZOS)
    buf = io.BytesIO()
    flat.save(buf, "PNG", optimize=True)
    return "data:image/png;base64," + base64.b64encode(buf.getvalue()).decode("ascii")


def encode_raw(path: Path, max_side: int = 720) -> str:
    im = Image.open(path).convert("RGB")
    im.thumbnail((max_side, max_side), Image.Resampling.LANCZOS)
    buf = io.BytesIO()
    im.save(buf, "JPEG", quality=90, optimize=True)
    return "data:image/jpeg;base64," + base64.b64encode(buf.getvalue()).decode("ascii")


def set_run_font(run, name="Calibri", size=14, bold=False, color=INK, italic=False):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_rect(slide, l, t, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def add_round(slide, l, t, w, h, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    # tighter corners
    try:
        sh.adjustments[0] = 0.08
    except Exception:
        pass
    return sh


def add_oval(slide, l, t, w, h, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def add_text(slide, l, t, w, h, text, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Calibri", italic=False, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}.get(anchor, "t"))
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_font(run, font, size, bold, color, italic)
    return box


def add_para(tf, text, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT, space_before=0, space_after=4, font="Calibri", italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    set_run_font(run, font, size, bold, color, italic)
    return p


def disable_shadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def prepare_exports():
    out = {}

    # Lockup and mark are already tight.
    for name, src in {
        "lockup.png": ASSETS / "vytra_logo_lockup.png",
        "mark.png": ASSETS / "vytra_mark.png",
    }.items():
        dest = EXPORTS / name
        save_display(src, dest, max_side=1600, pad=28)
        out[name] = dest

    # Official stacked art has a large white square. Crop to content.
    dest = EXPORTS / "official.png"
    save_display(ASSETS / "vytra_logo_official.png", dest, max_side=1600, pad=48)
    out["official.png"] = dest

    # App icons must stay square.
    for name, src in {
        "icon_light.png": ASSETS / "vytra_app_icon_light.png",
        "icon_dark.png": ASSETS / "vytra_app_icon_dark.png",
    }.items():
        dest = EXPORTS / name
        im = Image.open(src).convert("RGB")
        im.thumbnail((1200, 1200), Image.Resampling.LANCZOS)
        im.save(dest, "PNG", optimize=True)
        out[name] = dest
    return out


def _font(size, bold=False):
    from PIL import ImageFont
    path = (
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"
        if bold
        else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"
    )
    return ImageFont.truetype(path, size)


def _fit(im: Image.Image, box: tuple[int, int], bg=(255, 255, 255)) -> Image.Image:
    canvas = Image.new("RGB", box, bg)
    copy = im.convert("RGB")
    copy.thumbnail((box[0], box[1]), Image.Resampling.LANCZOS)
    x = (box[0] - copy.width) // 2
    y = (box[1] - copy.height) // 2
    canvas.paste(copy, (x, y))
    return canvas


def build_board(imgs: dict[str, Path]) -> Path:
    """Single poster that arranges every official asset."""
    from PIL import ImageDraw

    W, H = 3200, 2000
    FOREST_C = (14, 42, 28)
    DEEP = (8, 28, 18)
    PULSE_C = (108, 165, 50)
    CREAM_C = (244, 241, 234)
    WHITE_C = (255, 255, 255)
    SOFT_C = (232, 243, 220)
    MUTED_C = (71, 84, 103)
    INK_C = (15, 23, 42)

    img = Image.new("RGB", (W, H), CREAM_C)
    d = ImageDraw.Draw(img)

    # Left identity panel
    d.rectangle((0, 0, 1180, H), fill=DEEP)
    d.rectangle((0, 0, 18, H), fill=PULSE_C)
    d.ellipse((1080, 70, 1110, 100), fill=PULSE_C)

    d.text((80, 80), "BRAND IDENTITY  ·  1.0", font=_font(28, True), fill=PULSE_C)
    d.text((72, 180), "vytra", font=_font(148, True), fill=WHITE_C)
    d.text((80, 360), "See Health. Detect Early.", font=_font(42, False), fill=SOFT_C)

    d.text(
        (80, 460),
        "Offline smartphone screening aid\nfor ASHA workers. Not a diagnosis.",
        font=_font(30, False),
        fill=(181, 203, 178),
        spacing=10,
    )

    d.rounded_rectangle((80, 620, 1100, 760), radius=16, fill=FOREST_C)
    d.text((110, 655), "VYTRA  =  Vision + Vitality + Tracking + AI", font=_font(28, True), fill=WHITE_C)
    d.text((110, 700), "PITCH ONLY  ·  NEVER ON THE WORKER APP", font=_font(18, True), fill=PULSE_C)

    letters = [
        ("V", "Vision", "The camera is the sensor."),
        ("Y", "Your health", "Built for the person in the house."),
        ("T", "Tracking", "A dated note, not a verbal hunch."),
        ("R", "Recognition / Response", "A referral, not a prescription."),
        ("A", "AI / Analysis", "On-device colour maths."),
    ]
    y = 800
    for ch, word, blurb in letters:
        d.text((90, y), ch, font=_font(40, True), fill=PULSE_C)
        d.text((160, y + 4), word, font=_font(26, True), fill=WHITE_C)
        d.text((160, y + 42), blurb, font=_font(20, False), fill=(181, 203, 178))
        y += 92

    d.text((80, 1880), "SIH 2026  ·  TEAM VYTRA  ·  13 AUG 2026", font=_font(20, True), fill=(155, 179, 154))

    # Right: primary lockup plate
    d.rounded_rectangle((1240, 70, 3120, 1280), radius=24, fill=WHITE_C)
    d.text((1280, 100), "PRIMARY LOCKUP", font=_font(20, True), fill=MUTED_C)
    d.text((2480, 100), "vytra_logo_lockup.png", font=_font(20, False), fill=MUTED_C)
    lock = Image.open(imgs["lockup.png"]).convert("RGB")
    fitted = _fit(lock, (1700, 1040), WHITE_C)
    img.paste(fitted, (1330, 160))

    # Bottom row of three
    tiles = [
        (1240, "SYMBOL / MARK", "vytra_mark.png", imgs["mark.png"], WHITE_C),
        (1870, "APP ICON  ·  LIGHT", "vytra_app_icon_light.png", imgs["icon_light.png"], WHITE_C),
        (2500, "APP ICON  ·  DARK", "vytra_app_icon_dark.png", imgs["icon_dark.png"], FOREST_C),
    ]
    for x, title, fname, path, bg in tiles:
        d.rounded_rectangle((x, 1320, x + 620, 1920), radius=20, fill=bg)
        label_c = (183, 203, 180) if bg == FOREST_C else MUTED_C
        d.text((x + 28, 1345), title, font=_font(18, True), fill=label_c)
        inner = Image.open(path).convert("RGB")
        if "icon" in fname or "icon" in title.lower() or "app_icon" in str(path):
            fitted = _fit(inner, (420, 420), bg)
            img.paste(fitted, (x + 100, 1410))
        else:
            fitted = _fit(inner, (540, 400), bg)
            img.paste(fitted, (x + 40, 1420))

    dest = ROOT / "VYTRA_Brand_Board.png"
    img.save(dest, "PNG", optimize=True)
    return dest


def build_html(imgs: dict[str, Path]):
    lockup = encode_png(imgs["lockup.png"], 780, 20)
    mark = encode_png(imgs["mark.png"], 720, 16)
    icon_l = encode_raw(imgs["icon_light.png"], 640)
    icon_d = encode_raw(imgs["icon_dark.png"], 640)

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8" />
<meta name="viewport" content="width=device-width, initial-scale=1" />
<title>VYTRA — Brand Identity System</title>
<style>
  :root {{
    --forest: #0E2A1C;
    --forest-2: #081C12;
    --pulse: #6CA532;
    --pulse-2: #4F7E24;
    --soft: #E8F3DC;
    --cream: #F4F1EA;
    --paper: #FFFEFB;
    --ink: #0F172A;
    --muted: #475467;
    --line: #E4E7EC;
    --white: #ffffff;
    --warm: #EFEBE1;
  }}
  * {{ box-sizing: border-box; margin: 0; padding: 0; }}
  html {{ scroll-behavior: smooth; }}
  body {{
    font-family: "Segoe UI", "Helvetica Neue", Arial, sans-serif;
    background: var(--cream);
    color: var(--ink);
    line-height: 1.55;
    font-size: 16px;
  }}
  a {{ color: inherit; text-decoration: none; }}
  img {{ max-width: 100%; display: block; }}

  nav {{
    position: sticky; top: 0; z-index: 20;
    background: rgba(8,28,18,0.96);
    color: #fff;
    display: flex; align-items: center; justify-content: space-between;
    padding: 12px 40px;
    backdrop-filter: blur(8px);
  }}
  nav .brand {{
    display: flex; align-items: center; gap: 10px;
    letter-spacing: 0.28em; font-size: 12px; font-weight: 700;
  }}
  nav .dot {{ width: 8px; height: 8px; border-radius: 50%; background: var(--pulse); }}
  nav ul {{ display: flex; gap: 22px; list-style: none; font-size: 12px; letter-spacing: 0.08em; text-transform: uppercase; color: #cfe0c4; }}
  nav ul a:hover {{ color: #fff; }}

  .cover {{
    min-height: calc(100vh - 46px);
    background: radial-gradient(1200px 600px at 80% 10%, #163a27 0%, var(--forest) 42%, var(--forest-2) 100%);
    color: #fff;
    padding: 72px 72px 56px;
    display: grid;
    grid-template-columns: 1.1fr 0.9fr;
    gap: 48px;
    align-items: center;
  }}
  .kicker {{
    color: var(--pulse); letter-spacing: 0.34em; text-transform: uppercase;
    font-size: 12px; font-weight: 700; margin-bottom: 22px;
  }}
  .cover h1 {{
    font-size: 92px; line-height: 0.9; letter-spacing: -0.04em; font-weight: 700;
  }}
  .cover h1 span {{ color: var(--pulse); }}
  .cover .tag {{
    margin-top: 22px; font-size: 26px; font-weight: 500; color: #e8f3dc;
  }}
  .cover .lede {{
    margin-top: 18px; max-width: 520px; color: #c5d4c3; font-size: 17px;
  }}
  .meta-row {{
    margin-top: 48px; display: flex; gap: 28px; flex-wrap: wrap;
    font-size: 12px; letter-spacing: 0.12em; text-transform: uppercase; color: #9bb39a;
  }}
  .plate {{
    background: #fff; border-radius: 8px; padding: 48px 40px 40px;
    display: flex; align-items: center; justify-content: center;
    box-shadow: 0 30px 80px rgba(0,0,0,0.28);
  }}
  .plate img {{ width: 100%; max-width: 380px; height: auto; }}

  section.page {{
    padding: 88px 72px;
    max-width: 1280px;
    margin: 0 auto;
  }}
  .sec-label {{
    display: flex; align-items: baseline; gap: 16px; margin-bottom: 14px;
  }}
  .sec-num {{
    color: var(--pulse-2); font-weight: 700; letter-spacing: 0.2em; font-size: 12px;
  }}
  h2 {{
    font-size: 40px; letter-spacing: -0.03em; line-height: 1.1; color: var(--forest);
  }}
  .rule {{ height: 1px; background: #ddd6c8; margin: 22px 0 36px; }}
  .intro {{
    max-width: 720px; color: var(--muted); font-size: 18px; margin-bottom: 40px;
  }}

  .formula {{
    background: var(--forest);
    color: #fff;
    border-radius: 10px;
    padding: 28px 36px;
    display: flex; align-items: center; justify-content: space-between;
    gap: 20px; flex-wrap: wrap;
    margin-bottom: 28px;
  }}
  .formula .eq {{
    font-size: 22px; font-weight: 700; letter-spacing: 0.02em;
  }}
  .formula .eq b {{ color: var(--pulse); font-weight: 700; }}
  .formula .note {{ font-size: 12px; letter-spacing: 0.14em; text-transform: uppercase; color: #b7cbb4; }}

  .letters {{
    display: grid; grid-template-columns: repeat(5, 1fr); gap: 14px;
  }}
  .letter {{
    background: var(--paper);
    border: 1px solid #e7e1d4;
    border-radius: 10px;
    padding: 22px 18px 20px;
    min-height: 210px;
  }}
  .letter .ch {{
    font-size: 42px; font-weight: 700; color: var(--forest); letter-spacing: -0.04em;
  }}
  .letter .word {{
    margin-top: 10px; font-weight: 700; color: var(--forest); font-size: 15px;
  }}
  .letter p {{ margin-top: 8px; font-size: 13px; color: var(--muted); }}
  .letter .bar {{ width: 28px; height: 3px; background: var(--pulse); margin: 14px 0 0; }}

  .two {{ display: grid; grid-template-columns: 1fr 1fr; gap: 22px; }}
  .card {{
    background: var(--paper);
    border: 1px solid #e7e1d4;
    border-radius: 10px;
    padding: 28px;
  }}
  .card h3 {{ font-size: 13px; letter-spacing: 0.16em; text-transform: uppercase; color: var(--pulse-2); margin-bottom: 10px; }}
  .quote {{
    font-size: 34px; line-height: 1.15; letter-spacing: -0.03em; color: var(--forest); font-weight: 600;
  }}
  .card li {{ margin: 8px 0 0 18px; color: var(--muted); }}

  .family {{
    display: grid;
    grid-template-columns: 1.15fr 0.85fr;
    grid-template-rows: auto auto;
    gap: 16px;
  }}
  .tile {{
    background: #fff;
    border: 1px solid #e7e1d4;
    border-radius: 10px;
    padding: 22px;
    display: flex; flex-direction: column;
  }}
  .tile.span {{ grid-row: span 2; }}
  .tile .cap {{
    display: flex; justify-content: space-between; align-items: baseline;
    font-size: 11px; letter-spacing: 0.16em; text-transform: uppercase; color: var(--muted);
    margin-bottom: 18px;
  }}
  .tile .stage {{
    flex: 1; display: flex; align-items: center; justify-content: center;
    min-height: 220px;
  }}
  .tile .stage img {{ max-height: 360px; width: auto; }}
  .tile.dark {{ background: var(--forest); border-color: var(--forest); }}
  .tile.dark .cap {{ color: #b7cbb4; }}
  .icons {{ display: grid; grid-template-columns: 1fr 1fr; gap: 16px; margin-top: 16px; }}
  .icons .stage img {{ width: 72%; max-width: 260px; height: auto; }}

  .swatches {{ display: grid; grid-template-columns: repeat(6, 1fr); gap: 12px; }}
  .swatch {{ border-radius: 10px; overflow: hidden; border: 1px solid #e7e1d4; background: #fff; }}
  .swatch .chip {{ height: 92px; }}
  .swatch .meta {{ padding: 12px 12px 14px; }}
  .swatch .name {{ font-weight: 700; font-size: 13px; color: var(--forest); }}
  .swatch .hex {{ font-size: 12px; color: var(--muted); letter-spacing: 0.04em; margin-top: 2px; }}
  .swatch .role {{ font-size: 11px; color: #6b7280; margin-top: 6px; }}

  table {{
    width: 100%; border-collapse: collapse; background: #fff;
    border: 1px solid #e7e1d4; border-radius: 10px; overflow: hidden;
    font-size: 14px;
  }}
  th {{
    text-align: left; background: var(--forest); color: #fff;
    font-weight: 600; letter-spacing: 0.06em; font-size: 12px; text-transform: uppercase;
    padding: 12px 16px;
  }}
  td {{ padding: 12px 16px; border-top: 1px solid #eee8dc; vertical-align: top; }}
  tr:nth-child(even) td {{ background: #fbfaf6; }}
  .ok {{ color: #027A48; font-weight: 700; }}
  .no {{ color: #B42318; font-weight: 700; }}

  .donts {{ display: grid; grid-template-columns: 1fr 1fr; gap: 14px; }}
  .do, .dont {{
    border-radius: 10px; padding: 18px 20px; background: #fff; border: 1px solid #e7e1d4;
  }}
  .do {{ border-left: 4px solid #027A48; }}
  .dont {{ border-left: 4px solid #B42318; }}
  .do h4, .dont h4 {{ font-size: 13px; margin-bottom: 6px; }}
  .do p, .dont p {{ font-size: 14px; color: var(--muted); }}

  .anatomy {{
    display: grid; grid-template-columns: 1fr 1fr; gap: 22px; align-items: center;
  }}
  .callouts {{ display: flex; flex-direction: column; gap: 12px; }}
  .call {{
    display: grid; grid-template-columns: 36px 1fr; gap: 12px; align-items: start;
    background: #fff; border: 1px solid #e7e1d4; border-radius: 10px; padding: 14px 16px;
  }}
  .n {{
    width: 28px; height: 28px; border-radius: 50%; background: var(--forest); color: #fff;
    display: flex; align-items: center; justify-content: center; font-size: 12px; font-weight: 700;
  }}
  .call strong {{ display: block; color: var(--forest); font-size: 14px; }}
  .call span {{ font-size: 13px; color: var(--muted); }}

  footer.page {{
    background: var(--forest-2); color: #c5d4c3; padding: 48px 72px 56px;
  }}
  footer.page .row {{ display: flex; justify-content: space-between; gap: 24px; flex-wrap: wrap; }}
  footer.page strong {{ color: #fff; }}
  .disclaimer {{
    margin-top: 22px; font-size: 12px; color: #9bb39a; max-width: 820px;
  }}

  @media (max-width: 980px) {{
    nav {{ padding: 12px 18px; }}
    nav ul {{ display: none; }}
    .cover, section.page, footer.page {{ padding: 40px 20px; }}
    .cover {{ grid-template-columns: 1fr; min-height: auto; }}
    .cover h1 {{ font-size: 64px; }}
    .letters, .two, .family, .swatches, .donts, .anatomy, .icons {{ grid-template-columns: 1fr; }}
    .tile.span {{ grid-row: auto; }}
  }}
  @media print {{
    nav {{ position: static; }}
    .cover {{ min-height: auto; page-break-after: always; }}
    section.page {{ page-break-after: always; padding: 36px; }}
  }}
</style>
</head>
<body>
<nav>
  <div class="brand"><span class="dot"></span>VYTRA · IDENTITY</div>
  <ul>
    <li><a href="#meaning">Meaning</a></li>
    <li><a href="#logo">Logo</a></li>
    <li><a href="#colour">Colour</a></li>
    <li><a href="#type">Type</a></li>
    <li><a href="#use">Usage</a></li>
    <li><a href="#files">Files</a></li>
  </ul>
</nav>

<header class="cover">
  <div>
    <div class="kicker">Brand identity system · 01</div>
    <h1>vytra<span>.</span></h1>
    <div class="tag">See Health. Detect Early.</div>
    <p class="lede">Official lockup, mark, colour and naming for the VYTRA AI health-screening aid. Use these files. Do not redraw the V, the leaf-pulse, or the wordmark.</p>
    <div class="meta-row">
      <div>SIH 2026</div>
      <div>Team VYTRA</div>
      <div>Identity 1.0 · 13 Aug 2026</div>
      <div>Not a medical device</div>
    </div>
  </div>
  <div class="plate">
    <img src="{lockup}" alt="VYTRA official lockup" />
  </div>
</header>

<section class="page" id="meaning">
  <div class="sec-label"><span class="sec-num">01 / MEANING</span></div>
  <h2>What the name holds</h2>
  <div class="rule"></div>
  <p class="intro">VYTRA is a constructed name. The expansion is for the pitch deck and this book only — never inside the worker app, never on a result screen, never on the referral PDF.</p>

  <div class="formula">
    <div class="eq">VYTRA &nbsp;=&nbsp; <b>V</b>ision &nbsp;+&nbsp; <b>V</b>itality &nbsp;+&nbsp; <b>T</b>racking &nbsp;+&nbsp; <b>AI</b></div>
    <div class="note">Pitch only · not a clinical claim</div>
  </div>

  <div class="letters">
    <article class="letter">
      <div class="ch">V</div>
      <div class="word">Vision</div>
      <div class="bar"></div>
      <p>The camera is the sensor. We look at the eye, not a blood vial.</p>
    </article>
    <article class="letter">
      <div class="ch">Y</div>
      <div class="word">Your health</div>
      <div class="bar"></div>
      <p>Built for the person in the house. Used by the worker in front of them.</p>
    </article>
    <article class="letter">
      <div class="ch">T</div>
      <div class="word">Tracking</div>
      <div class="bar"></div>
      <p>Each visit becomes a dated, shareable note — not only a verbal hunch.</p>
    </article>
    <article class="letter">
      <div class="ch">R</div>
      <div class="word">Recognition / Response</div>
      <div class="bar"></div>
      <p>Recognise a colour signal. Respond with a referral, not a prescription.</p>
    </article>
    <article class="letter">
      <div class="ch">A</div>
      <div class="word">AI / Analysis</div>
      <div class="bar"></div>
      <p>On-device colour analysis. Not a cloud “AI doctor.”</p>
    </article>
  </div>
</section>

<section class="page" id="line" style="padding-top:0">
  <div class="two">
    <div class="card">
      <h3>Brand line</h3>
      <div class="quote">“See Health.<br>Detect Early.”</div>
      <p style="margin-top:16px;color:var(--muted)">Home screen, splash, title slide. Never on the result screen or the PDF — those surfaces carry a medical disclaimer, not a slogan.</p>
    </div>
    <div class="card">
      <h3>How the name is written</h3>
      <ul>
        <li><b>Logo / splash / home</b> — official PNG lockup. Never typeset a homemade V.</li>
        <li><b>SIH / college forms</b> — VYTRA (unique, no college in the name).</li>
        <li><b>Sentences in docs</b> — VYTRA</li>
        <li><b>Fallback text only</b> — vytra, Noto Sans Medium, lowercase</li>
        <li><b>Never</b> Vytra, VYTRA-AI, AnamoAI, or MRCET_VYTRA</li>
      </ul>
    </div>
  </div>
</section>

<section class="page" id="logo">
  <div class="sec-label"><span class="sec-num">02 / LOGO SYSTEM</span></div>
  <h2>One mark. Four files.</h2>
  <div class="rule"></div>
  <p class="intro">Forest-green V, lime leaf-pulse (heartbeat + leaf), lime living-point, lowercase wordmark, tracked category line. The designer note “vitality + diagnostics” is archive-only and must never ship.</p>

  <div class="family">
    <div class="tile span">
      <div class="cap"><span>Primary lockup</span><span>vytra_logo_lockup.png</span></div>
      <div class="stage"><img src="{lockup}" alt="Primary lockup" /></div>
    </div>
    <div class="tile">
      <div class="cap"><span>Symbol / mark</span><span>vytra_mark.png</span></div>
      <div class="stage"><img src="{mark}" alt="VYTRA mark" style="max-height:240px" /></div>
    </div>
    <div class="tile">
      <div class="cap"><span>Archive only — do not ship</span><span>vytra_logo_source.jpg</span></div>
      <div class="stage" style="min-height:160px;align-items:flex-start;padding-top:8px">
        <p style="color:var(--muted);font-size:14px;line-height:1.5;max-width:340px">The source file includes the designer note <i>“leaf-pulse concept — vitality + diagnostics.”</i> That line is not part of the logo. Never put it in the app, on a result, or on a slide.</p>
      </div>
    </div>
  </div>

  <div class="icons">
    <div class="tile">
      <div class="cap"><span>App icon · light</span><span>vytra_app_icon_light.png</span></div>
      <div class="stage"><img src="{icon_l}" alt="Light app icon" /></div>
    </div>
    <div class="tile dark">
      <div class="cap"><span>App icon · dark</span><span>vytra_app_icon_dark.png</span></div>
      <div class="stage"><img src="{icon_d}" alt="Dark app icon" /></div>
    </div>
  </div>
</section>

<section class="page" id="anatomy" style="padding-top:0">
  <div class="anatomy">
    <div class="tile">
      <div class="cap"><span>Construction</span><span>Do not separate the parts</span></div>
      <div class="stage"><img src="{mark}" alt="Mark anatomy" /></div>
    </div>
    <div class="callouts">
      <div class="call"><div class="n">1</div><div><strong>Forest V</strong><span>Primary letterform. Colour #0E2A1C. Never outline, bevel, or redraw.</span></div></div>
      <div class="call"><div class="n">2</div><div><strong>Leaf-pulse</strong><span>Heartbeat + leaf in #6CA532. Vitality and a clinical signal — not a diagnosis.</span></div></div>
      <div class="call"><div class="n">3</div><div><strong>Living point</strong><span>Lime dot at the upper right. Clear-space unit equals the height of this dot.</span></div></div>
      <div class="call"><div class="n">4</div><div><strong>Wordmark</strong><span>lowercase vytra in the same forest green. Do not set it in another font.</span></div></div>
      <div class="call"><div class="n">5</div><div><strong>Category line</strong><span>AI HEALTH SCREENING — title slide, splash, S01 only.</span></div></div>
    </div>
  </div>
</section>

<section class="page" id="colour">
  <div class="sec-label"><span class="sec-num">03 / COLOUR</span></div>
  <h2>Two greens are the identity</h2>
  <div class="rule"></div>
  <p class="intro">Sampled from the official artwork. Do not recolour the pulse to teal. Risk reds and ambers are clinical UI tokens, not brand colours — they never recolour the logo.</p>
  <div class="swatches">
    <div class="swatch"><div class="chip" style="background:#0E2A1C"></div><div class="meta"><div class="name">Forest</div><div class="hex">#0E2A1C</div><div class="role">V, wordmark, headers</div></div></div>
    <div class="swatch"><div class="chip" style="background:#6CA532"></div><div class="meta"><div class="name">Pulse</div><div class="hex">#6CA532</div><div class="role">Leaf-pulse, dot, accent</div></div></div>
    <div class="swatch"><div class="chip" style="background:#E8F3DC"></div><div class="meta"><div class="name">Soft</div><div class="hex">#E8F3DC</div><div class="role">Chips, selected state</div></div></div>
    <div class="swatch"><div class="chip" style="background:#FFFFFF;border-bottom:1px solid #eee"></div><div class="meta"><div class="name">Paper</div><div class="hex">#FFFFFF</div><div class="role">Logo always sits here</div></div></div>
    <div class="swatch"><div class="chip" style="background:#0F172A"></div><div class="meta"><div class="name">Ink</div><div class="hex">#0F172A</div><div class="role">Body text</div></div></div>
    <div class="swatch"><div class="chip" style="background:#475467"></div><div class="meta"><div class="name">Muted</div><div class="hex">#475467</div><div class="role">Secondary text</div></div></div>
  </div>
  <p style="margin-top:22px;color:var(--muted);font-size:14px">Risk tokens (UI only): High #B42318 · Moderate #B54708 · Low #027A48 · Unable #475467. Never apply these to the mark.</p>
</section>

<section class="page" id="type">
  <div class="sec-label"><span class="sec-num">04 / TYPE &amp; VOICE</span></div>
  <h2>Noto Sans. Calm, not clinical-chic.</h2>
  <div class="rule"></div>
  <div class="two">
    <div class="card">
      <h3>Product type</h3>
      <p style="font-size:42px;line-height:1.1;letter-spacing:-0.03em;color:var(--forest);font-weight:700">Noto Sans</p>
      <p style="margin-top:10px;color:var(--muted)">Regular + SemiBold, with Noto Sans Telugu for <b>te</b>. No Inter, Poppins, Montserrat, Roboto for Telugu.</p>
      <p style="margin-top:18px;letter-spacing:0.28em;font-weight:700;color:var(--forest)">VYTRA</p>
      <p style="margin-top:4px;color:var(--muted);font-size:13px">Fallback wordmark: uppercase, tracking +40, SemiBold. Prefer the PNG.</p>
    </div>
    <div class="card">
      <h3>Voice</h3>
      <ul>
        <li>Honest, rural, buildable.</li>
        <li>Screening aid — never diagnosis.</li>
        <li>Icons carry the verb; text confirms it.</li>
        <li>Banned on worker surfaces: diagnostics, clinically validated, 95% accurate, FDA / ICMR approved, hemoglobin numbers.</li>
        <li>Retired name AnamoAI is gone.</li>
      </ul>
    </div>
  </div>
</section>

<section class="page" id="use">
  <div class="sec-label"><span class="sec-num">05 / USAGE</span></div>
  <h2>Where each line may appear</h2>
  <div class="rule"></div>
  <table>
    <thead><tr><th>Line / file</th><th>Allowed</th><th>Forbidden</th></tr></thead>
    <tbody>
      <tr><td><b>Mark + vytra lockup</b></td><td>Splash, S01, S02, PPT title, PDF header (small)</td><td>On a photograph, on a risk tile</td></tr>
      <tr><td><b>AI HEALTH SCREENING</b></td><td>Under the lockup on S01 / title slide</td><td>Result screen, PDF body, consent</td></tr>
      <tr><td><i>See Health. Detect Early.</i></td><td>S01, S02, splash, under the lockup</td><td>Result screen, PDF</td></tr>
      <tr><td>Vision + Vitality + Tracking + AI</td><td>Pitch deck and this book</td><td>Worker app</td></tr>
      <tr><td>“leaf-pulse · vitality + diagnostics”</td><td>Designer archive only</td><td><span class="no">Nowhere in the product.</span> The word diagnostics is banned on worker surfaces.</td></tr>
      <tr><td>Dark icon</td><td>Dark splash only</td><td>Do not invert the two greens</td></tr>
    </tbody>
  </table>

  <div class="donts" style="margin-top:28px">
    <div class="do"><h4 class="ok">Do</h4><p>Keep clear space around the lockup equal to the height of the green dot. Sit the logo on white or near-white. Use the PNG. On dark grounds, place a white plate behind the lockup or use the dark app icon.</p></div>
    <div class="dont"><h4 class="no">Don’t</h4><p>Do not outline, bevel, drop-shadow, or put the mark on a busy photo. Do not separate the pulse from the V. Do not replace the pulse with a caduceus or a red cross. Do not recolour the pulse.</p></div>
    <div class="do"><h4 class="ok">App icon</h4><p>Launcher = mark on a white rounded square. Adaptive: foreground = mark, background #FFFFFF. Safe zone 66%. No wordmark on the small icon.</p></div>
    <div class="dont"><h4 class="no">Don’t invent a second mark</h4><p>No eye-only icon, no teal recolour, no “VYTRA AI” monogram. The locked artwork is the identity.</p></div>
  </div>
</section>

<section class="page" id="files">
  <div class="sec-label"><span class="sec-num">06 / FILE INDEX</span></div>
  <h2>Ship these. Archive the rest.</h2>
  <div class="rule"></div>
  <table>
    <thead><tr><th>File</th><th>Role</th><th>Use</th></tr></thead>
    <tbody>
      <tr><td><b>vytra_mark.png</b></td><td>Symbol</td><td>App icon, splash mark, slide corner, favicon</td></tr>
      <tr><td><b>vytra_logo_lockup.png</b></td><td>Primary lockup</td><td>Title slide, poster, first screen of the app</td></tr>
      <tr><td><b>vytra_logo_for_slides.png</b></td><td>Slide lockup</td><td>Identical to the primary lockup</td></tr>
      <tr><td>vytra_logo_official.png</td><td>Archive</td><td>Contains the designer note — <span class="no">do not ship</span></td></tr>
      <tr><td><b>vytra_app_icon_light.png</b></td><td>Launcher / light</td><td>Adaptive icon on light</td></tr>
      <tr><td><b>vytra_app_icon_dark.png</b></td><td>Dark splash</td><td>Dark splash only</td></tr>
      <tr><td>vytra_logo_source.jpg</td><td>Archive</td><td>Includes designer notes — <span class="no">do not ship</span></td></tr>
    </tbody>
  </table>
</section>

<footer class="page">
  <div class="row">
    <div>
      <strong>VYTRA</strong> · See Health. Detect Early.<br />
      Brand Identity 1.0 · 13 August 2026 · Team VYTRA · SIH 2026
    </div>
    <div style="text-align:right">
      Artwork locked 12 August 2026<br />
      Do not redraw the V, the pulse, or the wordmark.
    </div>
  </div>
  <p class="disclaimer">This screening result is not a medical diagnosis. It is a triage aid for trained health workers only. All results require confirmation by a qualified medical professional. Do not make treatment decisions based on this result alone. The designer phrase “vitality + diagnostics” is not a product claim.</p>
</footer>
</body>
</html>
"""
    dest = ROOT / "VYTRA_Brand_Identity.html"
    dest.write_text(html, encoding="utf-8")
    return dest


def build_pptx(imgs: dict[str, Path]):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    W = prs.slide_width
    H = prs.slide_height

    def slide():
        return prs.slides.add_slide(blank)

    def footer(s, page, dark=False):
        col = RGBColor(0x9B, 0xB3, 0x9A) if dark else MUTED
        add_text(s, Inches(0.6), Inches(7.12), Inches(8), Inches(0.28),
                 "VYTRA  ·  Brand Identity 1.0  ·  See Health. Detect Early.",
                 10, False, col)
        add_text(s, Inches(11.4), Inches(7.12), Inches(1.3), Inches(0.28),
                 f"{page:02d}", 10, True, col, PP_ALIGN.RIGHT)

    def kicker(s, text, l=0.6, t=0.32, dark=False):
        add_oval(s, Inches(l), Inches(t + 0.06), Inches(0.11), Inches(0.11), PULSE)
        add_text(s, Inches(l + 0.2), Inches(t), Inches(8), Inches(0.28),
                 text.upper(), 11, True, PULSE if not dark else PULSE, PP_ALIGN.LEFT)

    # ---------- 01 COVER ----------
    s = slide()
    add_rect(s, 0, 0, W, H, FOREST_DEEP)
    add_rect(s, 0, 0, Inches(0.14), H, PULSE)
    add_oval(s, Inches(12.55), Inches(0.42), Inches(0.22), Inches(0.22), PULSE)
    kicker(s, "Brand identity system  ·  01", 0.7, 0.42, True)
    add_text(s, Inches(0.7), Inches(1.35), Inches(7.2), Inches(1.6),
             "vytra", 84, True, WHITE)
    add_text(s, Inches(0.72), Inches(2.95), Inches(7), Inches(0.5),
             "See Health. Detect Early.", 26, False, SOFT)
    add_text(s, Inches(0.72), Inches(3.6), Inches(6.6), Inches(1.2),
             "Official lockup, mark, colour and naming.\nUse these files. Do not redraw the V, the leaf-pulse, or the wordmark.",
             16, False, RGBColor(0xC5, 0xD4, 0xC3))
    add_text(s, Inches(0.72), Inches(6.55), Inches(7.4), Inches(0.35),
             "SIH 2026   ·   Team VYTRA   ·   Identity 1.0  ·  13 Aug 2026   ·   Not a medical device",
             11, False, RGBColor(0x9B, 0xB3, 0x9A))
    # white plate with lockup
    add_round(s, Inches(8.15), Inches(1.15), Inches(4.55), Inches(5.15), WHITE)
    s.shapes.add_picture(str(imgs["lockup.png"]), Inches(8.55), Inches(1.55), Inches(3.75), Inches(4.34))

    # ---------- 02 MEANING ----------
    s = slide()
    add_rect(s, 0, 0, W, H, CREAM)
    kicker(s, "01  /  Meaning")
    add_text(s, Inches(0.6), Inches(0.62), Inches(12), Inches(0.55),
             "What the name holds", 32, True, FOREST)
    add_rect(s, Inches(0.6), Inches(1.25), Inches(12.1), Inches(0.015), LINE)

    add_round(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(1.05), FOREST)
    add_text(s, Inches(0.9), Inches(1.62), Inches(9.6), Inches(0.5),
             "VYTRA  =  Vision  +  Vitality  +  Tracking  +  AI", 22, True, WHITE)
    add_text(s, Inches(0.9), Inches(2.1), Inches(9.6), Inches(0.28),
             "PITCH DECK ONLY  ·  NEVER INSIDE THE WORKER APP, RESULT SCREEN, OR PDF", 10, True, SOFT)
    add_text(s, Inches(8.7), Inches(1.78), Inches(3.7), Inches(0.55),
             "Not a clinical claim", 12, True, PULSE, PP_ALIGN.RIGHT)

    letters = [
        ("V", "Vision", "The camera is the sensor. We look at the eye, not a blood vial."),
        ("Y", "Your health", "Built for the person in the house. Used by the worker in front of them."),
        ("T", "Tracking", "Each visit becomes a dated, shareable note — not only a verbal hunch."),
        ("R", "Recognition / Response", "Recognise a colour signal. Respond with a referral, not a prescription."),
        ("A", "AI / Analysis", "On-device colour analysis. Not a cloud “AI doctor.”"),
    ]
    x0 = 0.6
    gap = 0.16
    card_w = (12.1 - 4 * gap) / 5
    for i, (ch, word, blurb) in enumerate(letters):
        x = Inches(x0 + i * (card_w + gap))
        add_round(s, x, Inches(2.8), Inches(card_w), Inches(3.85), WHITE)
        add_text(s, x + Inches(0.18), Inches(2.95), Inches(card_w - 0.3), Inches(0.7),
                 ch, 36, True, FOREST)
        add_rect(s, x + Inches(0.2), Inches(3.7), Inches(0.32), Inches(0.045), PULSE)
        add_text(s, x + Inches(0.18), Inches(3.88), Inches(card_w - 0.32), Inches(0.7),
                 word, 14, True, FOREST)
        add_text(s, x + Inches(0.18), Inches(4.6), Inches(card_w - 0.32), Inches(1.8),
                 blurb, 12, False, MUTED)
    footer(s, 2)

    # ---------- 03 LINE + NAMING ----------
    s = slide()
    add_rect(s, 0, 0, W, H, CREAM)
    kicker(s, "01  /  Meaning")
    add_text(s, Inches(0.6), Inches(0.62), Inches(12), Inches(0.55),
             "Brand line and naming", 32, True, FOREST)
    add_rect(s, Inches(0.6), Inches(1.25), Inches(12.1), Inches(0.015), LINE)

    add_round(s, Inches(0.6), Inches(1.55), Inches(6.2), Inches(5.15), FOREST)
    add_text(s, Inches(0.95), Inches(1.85), Inches(5.5), Inches(0.3),
             "BRAND LINE", 11, True, PULSE)
    add_text(s, Inches(0.95), Inches(2.35), Inches(5.5), Inches(1.8),
             "“See Health.\nDetect Early.”", 36, True, WHITE)
    add_text(s, Inches(0.95), Inches(4.5), Inches(5.5), Inches(1.7),
             "Home, splash, title slide only.\nNever on the result screen or the referral PDF — those surfaces carry a medical disclaimer, not a slogan.",
             15, False, SOFT)

    add_round(s, Inches(7.05), Inches(1.55), Inches(5.65), Inches(5.15), WHITE)
    add_text(s, Inches(7.35), Inches(1.8), Inches(5.2), Inches(0.3),
             "HOW THE NAME IS WRITTEN", 11, True, PULSE_DEEP)
    rows = [
        ("Logo / splash / home", "Official PNG lockup. Never a homemade V."),
        ("SIH / college forms", "VYTRA — unique, no college in the name."),
        ("Sentences in docs", "VYTRA"),
        ("Fallback text only", "vytra  ·  Noto Sans Medium  ·  lowercase"),
        ("Never write", "Vytra   ·   VYTRA-AI   ·   AnamoAI   ·   MRCET_VYTRA"),
    ]
    y = 2.25
    for title, body in rows:
        add_text(s, Inches(7.35), Inches(y), Inches(5.15), Inches(0.28), title, 14, True, FOREST)
        add_text(s, Inches(7.35), Inches(y + 0.26), Inches(5.15), Inches(0.36), body, 13, False, MUTED)
        y += 0.78
    footer(s, 3)

    # ---------- 04 LOGO FAMILY ----------
    s = slide()
    add_rect(s, 0, 0, W, H, CREAM)
    kicker(s, "02  /  Logo system")
    add_text(s, Inches(0.6), Inches(0.62), Inches(12), Inches(0.5),
             "One mark. Four files.", 32, True, FOREST)
    add_rect(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.015), LINE)

    # four tiles
    tiles = [
        (0.6, "PRIMARY LOCKUP", "vytra_logo_lockup.png", imgs["lockup.png"], False, 4.55, 5.15),
        (5.3, "SYMBOL / MARK", "vytra_mark.png", imgs["mark.png"], False, 3.85, 5.15),
        (9.3, "APP ICON · LIGHT", "vytra_app_icon_light.png", imgs["icon_light.png"], False, 1.85, 2.48),
    ]
    # left lockup
    add_round(s, Inches(0.6), Inches(1.45), Inches(4.55), Inches(5.25), WHITE)
    add_text(s, Inches(0.82), Inches(1.58), Inches(4.1), Inches(0.28), "PRIMARY LOCKUP", 10, True, MUTED)
    add_text(s, Inches(0.82), Inches(6.32), Inches(4.1), Inches(0.22), "vytra_logo_lockup.png", 10, False, MUTED)
    s.shapes.add_picture(str(imgs["lockup.png"]), Inches(1.05), Inches(2.0), Inches(3.65), Inches(4.15))

    add_round(s, Inches(5.3), Inches(1.45), Inches(3.85), Inches(5.25), WHITE)
    add_text(s, Inches(5.5), Inches(1.58), Inches(3.5), Inches(0.28), "SYMBOL / MARK", 10, True, MUTED)
    add_text(s, Inches(5.5), Inches(6.32), Inches(3.5), Inches(0.22), "vytra_mark.png", 10, False, MUTED)
    s.shapes.add_picture(str(imgs["mark.png"]), Inches(5.55), Inches(2.55), Inches(3.35), Inches(2.49))

    add_round(s, Inches(9.3), Inches(1.45), Inches(3.4), Inches(2.52), WHITE)
    add_text(s, Inches(9.48), Inches(1.54), Inches(3.05), Inches(0.22), "APP ICON · LIGHT", 10, True, MUTED)
    s.shapes.add_picture(str(imgs["icon_light.png"]), Inches(10.15), Inches(1.85), Inches(1.7), Inches(1.7))

    add_round(s, Inches(9.3), Inches(4.15), Inches(3.4), Inches(2.55), FOREST)
    add_text(s, Inches(9.48), Inches(4.24), Inches(3.05), Inches(0.22), "APP ICON · DARK", 10, True, SOFT)
    s.shapes.add_picture(str(imgs["icon_dark.png"]), Inches(10.15), Inches(4.55), Inches(1.7), Inches(1.7))
    footer(s, 4)

    # ---------- 05 LOCKUP LARGE ----------
    s = slide()
    add_rect(s, 0, 0, W, H, CREAM)
    kicker(s, "02  /  Logo system")
    add_text(s, Inches(0.6), Inches(0.62), Inches(12), Inches(0.5),
             "Primary lockup", 32, True, FOREST)
    add_rect(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.015), LINE)

    add_round(s, Inches(0.6), Inches(1.45), Inches(7.5), Inches(5.25), WHITE)
    s.shapes.add_picture(str(imgs["lockup.png"]), Inches(1.55), Inches(1.7), Inches(5.55), Inches(4.75))

    add_round(s, Inches(8.3), Inches(1.45), Inches(4.4), Inches(5.25), WHITE)
    add_text(s, Inches(8.55), Inches(1.7), Inches(4.0), Inches(0.28), "USE ON", 11, True, PULSE_DEEP)
    for i, line in enumerate([
        "Title slide of the SIH deck",
        "App splash and S01 / S02",
        "Poster and first screen",
        "PDF header — small only",
    ]):
        add_text(s, Inches(8.55), Inches(2.1 + i * 0.38), Inches(4.0), Inches(0.36), "▸  " + line, 14, False, INK)

    add_text(s, Inches(8.55), Inches(3.75), Inches(4.0), Inches(0.28), "NEVER ON", 11, True, RGBColor(0xB4, 0x23, 0x18))
    for i, line in enumerate([
        "A photograph or busy image",
        "A risk tile or result card",
        "Consent or PDF body",
        "Next to a homemade wordmark",
    ]):
        add_text(s, Inches(8.55), Inches(4.15 + i * 0.38), Inches(4.0), Inches(0.36), "▸  " + line, 14, False, INK)

    add_text(s, Inches(8.55), Inches(5.85), Inches(4.0), Inches(0.55),
             "Category line AI HEALTH SCREENING travels with this lockup only.",
             12, False, MUTED)
    footer(s, 5)

    # ---------- 06 ANATOMY ----------
    s = slide()
    add_rect(s, 0, 0, W, H, CREAM)
    kicker(s, "02  /  Logo system")
    add_text(s, Inches(0.6), Inches(0.62), Inches(12), Inches(0.5),
             "Anatomy — do not separate the parts", 32, True, FOREST)
    add_rect(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.015), LINE)

    add_round(s, Inches(0.6), Inches(1.5), Inches(6.5), Inches(5.15), WHITE)
    s.shapes.add_picture(str(imgs["mark.png"]), Inches(1.15), Inches(2.15), Inches(5.4), Inches(4.01))

    parts = [
        ("01", "Forest V", "Primary letterform. #0E2A1C. Never outline, bevel, or redraw."),
        ("02", "Leaf-pulse", "Heartbeat + leaf in #6CA532. A clinical signal — not a diagnosis."),
        ("03", "Living point", "Lime dot, upper right. Clear space = the height of this dot."),
        ("04", "Wordmark", "lowercase vytra, same forest green. Do not set it in another font."),
        ("05", "Category line", "AI HEALTH SCREENING — title slide, splash, S01 only."),
    ]
    y = 1.5
    for num, title, body in parts:
        add_round(s, Inches(7.35), Inches(y), Inches(5.35), Inches(0.95), WHITE)
        add_oval(s, Inches(7.52), Inches(y + 0.3), Inches(0.36), Inches(0.36), FOREST)
        add_text(s, Inches(7.52), Inches(y + 0.34), Inches(0.36), Inches(0.3), num[-1], 11, True, WHITE, PP_ALIGN.CENTER)
        add_text(s, Inches(8.05), Inches(y + 0.12), Inches(4.45), Inches(0.3), title, 14, True, FOREST)
        add_text(s, Inches(8.05), Inches(y + 0.44), Inches(4.45), Inches(0.42), body, 12, False, MUTED)
        y += 1.04
    footer(s, 6)

    # ---------- 07 COLOUR ----------
    s = slide()
    add_rect(s, 0, 0, W, H, CREAM)
    kicker(s, "03  /  Colour")
    add_text(s, Inches(0.6), Inches(0.62), Inches(12), Inches(0.5),
             "Two greens are the identity", 32, True, FOREST)
    add_rect(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.015), LINE)

    swatches = [
        ("Forest", "#0E2A1C", "V, wordmark, headers", RGBColor(0x0E, 0x2A, 0x1C), WHITE),
        ("Pulse", "#6CA532", "Leaf-pulse, dot, accent", RGBColor(0x6C, 0xA5, 0x32), WHITE),
        ("Soft", "#E8F3DC", "Chips, selected state", RGBColor(0xE8, 0xF3, 0xDC), FOREST),
        ("Paper", "#FFFFFF", "Logo always sits here", WHITE, FOREST),
        ("Ink", "#0F172A", "Body text", RGBColor(0x0F, 0x17, 0x2A), WHITE),
        ("Muted", "#475467", "Secondary text", RGBColor(0x47, 0x54, 0x67), WHITE),
    ]
    x = 0.6
    for name, hexv, role, fill, tc in swatches:
        add_round(s, Inches(x), Inches(1.55), Inches(1.95), Inches(3.55), WHITE)
        # color block
        sh = add_rect(s, Inches(x + 0.12), Inches(1.7), Inches(1.71), Inches(1.85), fill)
        if name == "Paper":
            sh.line.color.rgb = LINE
            sh.line.width = Pt(1)
        add_text(s, Inches(x + 0.16), Inches(3.7), Inches(1.65), Inches(0.3), name, 14, True, FOREST)
        add_text(s, Inches(x + 0.16), Inches(4.0), Inches(1.65), Inches(0.25), hexv, 12, False, MUTED)
        add_text(s, Inches(x + 0.16), Inches(4.3), Inches(1.65), Inches(0.6), role, 11, False, MUTED)
        x += 2.1

    add_text(s, Inches(0.6), Inches(5.35), Inches(12.1), Inches(0.35),
             "Do not recolour the pulse to teal. The two greens are the identity.",
             14, True, FOREST)
    add_text(s, Inches(0.6), Inches(5.75), Inches(12.1), Inches(0.7),
             "Risk tokens are clinical UI only — never apply them to the mark.\nHigh  #B42318     Moderate  #B54708     Low  #027A48     Unable  #475467",
             13, False, MUTED)
    footer(s, 7)

    # ---------- 08 USAGE MATRIX ----------
    s = slide()
    add_rect(s, 0, 0, W, H, CREAM)
    kicker(s, "04  /  Usage")
    add_text(s, Inches(0.6), Inches(0.62), Inches(12), Inches(0.5),
             "Where each line may appear", 32, True, FOREST)
    add_rect(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.015), LINE)

    # table header
    add_rect(s, Inches(0.6), Inches(1.45), Inches(12.1), Inches(0.46), FOREST)
    add_text(s, Inches(0.8), Inches(1.54), Inches(3.4), Inches(0.3), "LINE / FILE", 11, True, WHITE)
    add_text(s, Inches(4.3), Inches(1.54), Inches(4.2), Inches(0.3), "ALLOWED", 11, True, WHITE)
    add_text(s, Inches(8.6), Inches(1.54), Inches(3.9), Inches(0.3), "FORBIDDEN", 11, True, WHITE)

    rows = [
        ("Mark + vytra lockup", "Splash, S01, S02, PPT title, PDF header (small)", "On a photograph or a risk tile"),
        ("AI HEALTH SCREENING", "Under the lockup on S01 / title slide", "Result screen, PDF body, consent"),
        ("See Health. Detect Early.", "S01, S02, splash, under the lockup", "Result screen, PDF"),
        ("Vision + Vitality + Tracking + AI", "Pitch deck and this book", "Worker app"),
        ("“vitality + diagnostics”", "Designer archive only", "Nowhere in the product"),
        ("Dark app icon", "Dark splash only", "Do not invert the two greens"),
    ]
    y = 1.91
    for i, (a, b, c) in enumerate(rows):
        bg = WHITE if i % 2 == 0 else WARM
        add_rect(s, Inches(0.6), Inches(y), Inches(12.1), Inches(0.72), bg)
        add_text(s, Inches(0.8), Inches(y + 0.18), Inches(3.4), Inches(0.42), a, 13, True, FOREST)
        add_text(s, Inches(4.3), Inches(y + 0.18), Inches(4.2), Inches(0.42), b, 13, False, INK)
        add_text(s, Inches(8.6), Inches(y + 0.18), Inches(3.9), Inches(0.42), c, 13, False, RGBColor(0xB4, 0x23, 0x18))
        y += 0.72
    footer(s, 8)

    # ---------- 09 DO / DON'T ----------
    s = slide()
    add_rect(s, 0, 0, W, H, CREAM)
    kicker(s, "04  /  Usage")
    add_text(s, Inches(0.6), Inches(0.62), Inches(12), Inches(0.5),
             "Clear space, icons, mistakes", 32, True, FOREST)
    add_rect(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.015), LINE)

    add_round(s, Inches(0.6), Inches(1.5), Inches(6.05), Inches(5.15), WHITE)
    add_rect(s, Inches(0.6), Inches(1.5), Inches(0.12), Inches(5.15), RGBColor(0x02, 0x7A, 0x48))
    add_text(s, Inches(1.0), Inches(1.7), Inches(5.4), Inches(0.35), "DO", 16, True, RGBColor(0x02, 0x7A, 0x48))
    dos = [
        "Keep clear space equal to the height of the green dot.",
        "Sit the logo on white or very near-white.",
        "Use the official PNG. Do not redraw.",
        "On dark grounds: white plate, or the dark app icon.",
        "Launcher = mark on a white rounded square.",
        "Adaptive icon: foreground = mark, background #FFFFFF. Safe zone 66%.",
        "No wordmark on the small app icon — it will collapse.",
    ]
    yy = 2.2
    for d in dos:
        add_text(s, Inches(1.0), Inches(yy), Inches(5.35), Inches(0.5), "●   " + d, 14, False, INK)
        yy += 0.55

    add_round(s, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.15), WHITE)
    add_rect(s, Inches(6.9), Inches(1.5), Inches(0.12), Inches(5.15), RGBColor(0xB4, 0x23, 0x18))
    add_text(s, Inches(7.3), Inches(1.7), Inches(5.15), Inches(0.35), "DON’T", 16, True, RGBColor(0xB4, 0x23, 0x18))
    donts = [
        "Do not outline, bevel, or add a drop shadow.",
        "Do not place the mark on a busy photograph.",
        "Do not separate the pulse from the V.",
        "Do not replace the pulse with a caduceus or red cross.",
        "Do not invert or recolour the two greens.",
        "Do not invent a second mark or an eye-only icon.",
        "Do not print “diagnostics” on any worker surface.",
    ]
    yy = 2.2
    for d in donts:
        add_text(s, Inches(7.3), Inches(yy), Inches(5.15), Inches(0.5), "●   " + d, 14, False, INK)
        yy += 0.55
    footer(s, 9)

    # ---------- 10 FILE INDEX ----------
    s = slide()
    add_rect(s, 0, 0, W, H, CREAM)
    kicker(s, "05  /  Files")
    add_text(s, Inches(0.6), Inches(0.62), Inches(12), Inches(0.5),
             "Ship these. Archive the rest.", 32, True, FOREST)
    add_rect(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.015), LINE)

    add_rect(s, Inches(0.6), Inches(1.45), Inches(12.1), Inches(0.46), FOREST)
    add_text(s, Inches(0.8), Inches(1.54), Inches(4.2), Inches(0.3), "FILE", 11, True, WHITE)
    add_text(s, Inches(5.1), Inches(1.54), Inches(2.4), Inches(0.3), "ROLE", 11, True, WHITE)
    add_text(s, Inches(7.6), Inches(1.54), Inches(4.8), Inches(0.3), "USE", 11, True, WHITE)

    files = [
        ("vytra_mark.png", "Symbol", "App icon, splash mark, slide corner, favicon"),
        ("vytra_logo_lockup.png", "Primary lockup", "Title slide, poster, first screen of the app"),
        ("vytra_logo_official.png", "Stacked official", "Square crop of the same lockup"),
        ("vytra_logo_for_slides.png", "Slide lockup", "Identical to the primary lockup"),
        ("vytra_app_icon_light.png", "Launcher / light", "Adaptive / launcher on light"),
        ("vytra_app_icon_dark.png", "Dark splash", "Dark splash only"),
        ("vytra_logo_source.jpg", "Archive only", "Designer notes — do not ship"),
    ]
    y = 1.91
    for i, (a, b, c) in enumerate(files):
        bg = WHITE if i % 2 == 0 else WARM
        add_rect(s, Inches(0.6), Inches(y), Inches(12.1), Inches(0.62), bg)
        add_text(s, Inches(0.8), Inches(y + 0.15), Inches(4.2), Inches(0.36), a, 13, True, FOREST)
        add_text(s, Inches(5.1), Inches(y + 0.15), Inches(2.4), Inches(0.36), b, 13, False, INK)
        col = RGBColor(0xB4, 0x23, 0x18) if "do not" in c.lower() else INK
        add_text(s, Inches(7.6), Inches(y + 0.15), Inches(4.8), Inches(0.36), c, 13, False, col)
        y += 0.62

    add_text(s, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.45),
             "Artwork locked 12 August 2026. Do not redraw the V, the pulse, or the wordmark in another font.",
             13, False, MUTED)
    footer(s, 10)

    # ---------- 11 CLOSING ----------
    s = slide()
    add_rect(s, 0, 0, W, H, FOREST_DEEP)
    add_rect(s, 0, 0, Inches(0.14), H, PULSE)
    add_oval(s, Inches(12.55), Inches(0.42), Inches(0.22), Inches(0.22), PULSE)
    add_text(s, Inches(0.75), Inches(1.7), Inches(12), Inches(0.4),
             "TEAM VYTRA  ·  SIH 2026", 13, True, PULSE)
    add_text(s, Inches(0.72), Inches(2.2), Inches(12), Inches(1.4),
             "See Health. Detect Early.", 40, True, WHITE)
    add_text(s, Inches(0.75), Inches(3.7), Inches(11.2), Inches(1.3),
             "Offline smartphone screening aid for ASHA workers.\nNot a diagnosis. Not a medical device.",
             18, False, SOFT)
    add_text(s, Inches(0.75), Inches(5.4), Inches(11.5), Inches(1.1),
             "This screening result is not a medical diagnosis. It is a triage aid for trained health workers only.\nAll results require confirmation by a qualified medical professional.",
             13, False, RGBColor(0x9B, 0xB3, 0x9A))
    add_text(s, Inches(0.75), Inches(6.7), Inches(11.5), Inches(0.3),
             "Identity 1.0  ·  13 August 2026  ·  Artwork locked  ·  Do not redraw",
             12, False, RGBColor(0x7A, 0x96, 0x78))

    out = ROOT / "VYTRA_Brand_Identity.pptx"
    prs.save(out)
    return out


def main():
    imgs = prepare_exports()
    html = build_html(imgs)
    pptx = build_pptx(imgs)
    board = build_board(imgs)
    print("HTML", html, html.stat().st_size)
    print("PPTX", pptx, pptx.stat().st_size)
    print("BOARD", board, board.stat().st_size)
    for k, v in imgs.items():
        print(" ", k, Image.open(v).size, v.stat().st_size)


if __name__ == "__main__":
    main()
